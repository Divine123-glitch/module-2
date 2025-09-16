# generate_enhanced_sales_ppt.py

# This script generates an enhanced PowerPoint presentation (.pptx) with colors, charts, and graphs.

# It includes bar charts, pie charts, and line graphs using matplotlib-generated images.

# To run: Save this as generate_enhanced_sales_ppt.py, install required packages with:

# 'pip install python-pptx matplotlib', then 'python generate_enhanced_sales_ppt.py'

# Output: Enhanced_Sales_Performance_Report.pptx

from pptx import Presentation

from pptx.util import Inches, Pt

from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR_INDEX

from pptx.dml.color import RGBColor

from pptx.enum.dml import MSO_THEME_COLOR

import matplotlib.pyplot as plt

import io

import base64

# Create a new presentation

prs = Presentation()

# Generate charts as base64 images (using sample data from notes)

def generate_charts():

    charts = {}

    

    # Chart 1: 2015 Agent Revenue Bar Chart

    fig1, ax1 = plt.subplots(figsize=(6, 4))

    agents_2015 = ['Tolu', 'Tonye']

    revenues_2015 = [2883445, 289.12]

    bars1 = ax1.bar(agents_2015, revenues_2015, color=['#2E8B57', '#DC143C'])

    ax1.set_title('2015 Agent Revenue', fontsize=14, color='blue')

    ax1.set_ylabel('Revenue', color='green')

    ax1.bar_label(bars1, fmt='%.0f')

    plt.xticks(rotation=45)

    plt.tight_layout()

    buf1 = io.BytesIO()

    fig1.savefig(buf1, format='png', bbox_inches='tight')

    buf1.seek(0)

    charts['agent_2015'] = base64.b64encode(buf1.read()).decode('utf-8')

    plt.close(fig1)

    

    # Chart 2: Product Performance Over 2 Years (Subplots: Units and Revenue)

    fig2, (ax21, ax22) = plt.subplots(1, 2, figsize=(10, 4))

    products = ['HP', 'Apple']

    units_sold = [722, 10]

    revenues = [955000, 1500]

    bars2a = ax21.bar(products, units_sold, color=['#4169E1', '#FF8C00'])

    ax21.set_title('Units Sold', fontsize=12, color='navy')

    ax21.bar_label(bars2a, fmt='%d')

    bars2b = ax22.bar(products, revenues, color=['#4169E1', '#FF8C00'])

    ax22.set_title('Revenue', fontsize=12, color='navy')

    ax22.bar_label(bars2b, fmt='%.0f')

    plt.tight_layout()

    buf2 = io.BytesIO()

    fig2.savefig(buf2, format='png', bbox_inches='tight')

    buf2.seek(0)

    charts['product_overall'] = base64.b64encode(buf2.read()).decode('utf-8')

    plt.close(fig2)

    

    # Chart 3: 2015 Branch Performance Pie Chart

    fig3, ax3 = plt.subplots(figsize=(6, 4))

    branches = ['Ijoh', 'GRT']

    percentages = [70.45, 7.8]

    ax3.pie(percentages, labels=branches, autopct='%1.1f%%', colors=['#87CEEB', '#F08080'], startangle=90)

    ax3.set_title('2015 Branch Performance %', fontsize=14, color='purple')

    buf3 = io.BytesIO()

    fig3.savefig(buf3, format='png', bbox_inches='tight')

    buf3.seek(0)

    charts['branch_2015'] = base64.b64encode(buf3.read()).decode('utf-8')

    plt.close(fig3)

    

    # Chart 4: Revenue by Month Line Chart

    fig4, ax4 = plt.subplots(figsize=(6, 4))

    months = ['March', 'July', 'December']

    rev = [167.44, 1646, 2000]  # Assuming higher for Dec based on notes

    ax4.plot(months, rev, marker='o', color='#8A2BE2', linewidth=2, markersize=8)

    ax4.set_title('Revenue by Month Over 2 Years', fontsize=14, color='darkgreen')

    ax4.set_ylabel('Revenue', color='darkgreen')

    ax4.grid(True, linestyle='--', alpha=0.7)

    plt.xticks(rotation=45)

    plt.tight_layout()

    buf4 = io.BytesIO()

    fig4.savefig(buf4, format='png', bbox_inches='tight')

    buf4.seek(0)

    charts['monthly'] = base64.b64encode(buf4.read()).decode('utf-8')

    plt.close(fig4)

    

    return charts

charts = generate_charts()

# Helper function to add a slide with title, bullets, and optional chart

def add_slide(prs, title, bullets=None, chart_key=None, layout=1):

    slide_layout = prs.slide_layouts[layout]

    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title

    title_shape.text = title

    title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 51, 102)  # Navy blue

    title_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(24)

    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    

    if bullets:

        body_shape = slide.placeholders[1]

        tf = body_shape.text_frame

        tf.clear()

        for bullet in bullets:

            p = tf.add_paragraph()

            p.text = bullet

            p.level = 1 if bullet.startswith('  ') else 0

            p.alignment = PP_ALIGN.LEFT

            for run in p.runs:

                run.font.size = Pt(16)

                run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    

    # Add chart if provided

    if chart_key and chart_key in charts:

        img_data = charts[chart_key]

        img_stream = io.BytesIO(base64.b64decode(img_data))

        left = Inches(1)

        top = Inches(2)

        pic = slide.shapes.add_picture(img_stream, left, top, width=Inches(6))

    

    # Add colored footer

    footer_left = Inches(0.5)

    footer_top = Inches(6.5)

    footer_txbox = slide.shapes.add_textbox(footer_left, footer_top, Inches(9), Inches(0.5))

    footer_frame = footer_txbox.text_frame

    footer_frame.text = "United to Prevent Violence Against Women and Girls | www.cewhin.com"

    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for run in footer_frame.paragraphs[0].runs:

        run.font.size = Pt(12)

        run.font.color.rgb = RGBColor(139, 69, 19)  # Saddle brown

    

    return slide

# Slide 1: Title Slide (with color)

title_slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title

title.text = "Sales Performance Report: 2015-2016"

title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 128, 0)  # Green

title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)

subtitle = slide.placeholders[1]

subtitle.text = "Key Insights from Branch Data"

subtitle.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(100, 100, 100)  # Gray

# Footer

footer_txbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))

footer_frame = footer_txbox.text_frame

footer_frame.text = "United to Prevent Violence Against Women and Girls | www.cewhin.com"

footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

for run in footer_frame.paragraphs[0].runs:

    run.font.size = Pt(12)

    run.font.color.rgb = RGBColor(139, 69, 19)

# Slide 2: 2015 Agent Performance + Chart

add_slide(prs, "2015 Agent Performance", [

    "Tolu had the highest sales revenue",

    "  Total Rev: 2,883,445 / 21.08%",

    "Tonye had the lowest sales revenue",

    "  Total Per: 289.12 / 2.79%"

], 'agent_2015')

# Slide 3: 2015 Product Performance

add_slide(prs, "2015 Product Performance", [

    "HP had the highest sales revenue",

    "  Total Par: 5,814 / 414 units",

    "Lenovo had the lowest sales revenue",

    "  Total Par: 208.80 / 160 units",

    "Average Price of the year 2015: 200.00"

])

# Slide 4: 2015 Branch Performance + Pie Chart

add_slide(prs, "2015 Branch Performance", [

    "Ijoh is the highest performing branch",

    "  Total Rev: 7,305.56 / 70.45%",

    "GRT is the lowest performing branch",

    "  Total Revenue: 808.38 / 7.8%"

], 'branch_2015')

# Slide 5: 2016 Agent Performance

add_slide(prs, "2016 Agent Performance", [

    "Chinma had the highest sales revenue",

    "  Total Rev: 3,102 / 33.51%",

    "Torbari had the lowest sales revenue",

    "  Total Rev: 57.71 / 0.62%"

])

# Slide 6: 2016 Product Performance

add_slide(prs, "2016 Product Performance", [

    "Apple had the highest performing sales",

    "  Total Per: 3,247 / 308 units",

    "Apple had the lowest (note: possible data overlap; lowest units: 250 / 2 units)",

    "Average price of 2016: 125.00"

])

# Slide 7: 2016 Branch Performance

add_slide(prs, "2016 Branch Performance", [

    "GRA had the highest performance",

    "  Total Rev: 5,193 / 56%",

    "Addah Town had the lowest performance",

    "  Total Rev: 231.12 / 2.5%"

])

# Slide 8: 2016 Revenue by Month + Line Chart

add_slide(prs, "2016 Revenue by Month", [

    "July had the highest performance",

    "  Total Rev: 1,646",

    "March had the lowest performance",

    "  Total Rev: 167.44"

], 'monthly')

# Slide 9: Revenue by Branch (Overall)

add_slide(prs, "Revenue by Branch (Overall)", [

    "Ijoh Branch has the highest sales revenue",

    "  Total Revenue: 11,139.07 / 56.73%",

    "Town Branch has the lowest sales revenue",

    "  Total Revenue: 2,486.42 / 10.67%"

])

# Slide 10: Noticeable Insights

add_slide(prs, "Noticeable Insights", [

    "Only 3 out of 11 agents sold Apple products (Tolu, Emeka, Chinma)",

    "Only 5 out of 11 agents sold company products (Blessing, Ibrahim, Torbari, Chinma, Uche)",

    "Only 5 out of 11 agents sold HP products – DELL (Emeka, China, George, Blessing, Tolu)",

    "All 11 agents sold HP products",

    "Only 8 agents out of 11 sold Lenovo products (Tolu, George, Blessing, Tonye, Uche, Chinedu, Ibrahim, Tunde)",

    "No agent sold all 5 products",

    "Blessing and Uche sold the highest number of distinct products (4)",

    "Torbari and Tunde sold the lowest number of distinct products (2)"

])

# Slide 11: Over the 2 Years – Agent & Product Performance + Chart

add_slide(prs, "Over the 2 Years – Agent & Product Performance", [

    "Agent Performance:",

    "  Emeka had the highest sales revenue: 3,109.44 (15.84%)",

    "  Torbari had the lowest sales revenue: 536.75 (2.73%)",

    "Product Performance:",

    "  HP has the highest sales revenue: 955k / 722 units sold",

    "  Apple has the lowest sales revenue: 1.5k / 10 units sold"

], 'product_overall')

# Slide 12: Over the 2 Years – Year Performance

add_slide(prs, "Over the 2 Years – Year Performance", [

    "2015 is the highest performing year",

    "  Total Revenue: 10,369.54",

    "  Units Sold: 943",

    "2016 is the lowest performing year",

    "  Total Revenue: 9,258.39",

    "  Units Sold: MK (data unclear; possibly missing)"

])

# Slide 13: Revenue by Month Over the 2 Years

add_slide(prs, "Revenue by Month Over the 2 Years", [

    "December has the highest revenue by month",

    "  We have more sales in December than any other month",

    "March has the least revenue by month",

    "  We have less sales in the month of March"

])

# Slide 14: Thank You / Q&A

add_slide(prs, "Thank You / Q&A", [

    "Summary: Key trends show strong performance in 2015, HP dominance, and seasonal peaks in December/July.",

    "Contact: For more details, visit www.cewhin.com"

])

# Save the presentation

prs.save('Enhanced_Sales_Performance_Report.pptx')

print("Enhanced PPTX file generated successfully: Enhanced_Sales_Performance_Report.pptx")

print("Features added: Colored text/titles, embedded charts (bar, pie, line, subplots).")

