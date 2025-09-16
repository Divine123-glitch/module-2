from pptx import Presentation

from pptx.util import Inches, Pt

from pptx.enum.shapes import MSO_SHAPE

from pptx.dml.color import RGBColor

# Create presentation

prs = Presentation()

# Define slide layout

title_slide_layout = prs.slide_layouts[0]

bullet_slide_layout = prs.slide_layouts[1]

# Title Slide

slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title

subtitle = slide.placeholders[1]

title.text = "Sales Performance Report: 2015-2016"

subtitle.text = "A Comprehensive Review of Agent, Product, and Branch Performance"

# Footer

for shape in slide.shapes:

    if not shape.has_text_frame:

        continue

    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:

        for run in paragraph.runs:

            run.font.size = Pt(24)

# Slide content definitions

slides_content = [

    ("Executive Summary", [

        "Overview of sales performance across 2015 and 2016",

        "Key growth metrics and highlights"

    ]),

    ("Agent Performance Overview", [

        "Top-performing agents by revenue",

        "Year-over-year comparisons",

        "Percentage contribution to total sales"

    ]),

    ("Agent Performance Comparison", [

        "Bar chart of agent sales in 2015 vs 2016",

        "Growth rates and performance shifts"

    ]),

    ("Product Performance Overview", [

        "Best-selling products",

        "Revenue contribution by product category"

    ]),

    ("Product Performance Comparison", [

        "Pie charts or bar graphs for 2015 vs 2016",

        "Percentage changes and trends"

    ]),

    ("Branch Performance Overview", [

        "Sales by branch",

        "Regional performance highlights"

    ]),

    ("Branch Performance Comparison", [

        "Comparative analysis of branch growth",

        "Top and bottom performers"

    ]),

    ("Monthly Sales Trends: 2015", [

        "Line chart of monthly sales",

        "Seasonal patterns and anomalies"

    ]),

    ("Monthly Sales Trends: 2016", [

        "Line chart of monthly sales",

        "Notable changes from 2015"

    ]),

    ("Insights and Observations", [

        "Key drivers of performance",

        "Challenges and opportunities"

    ]),

    ("2-Year Summary", [

        "Total sales, growth rate, average monthly sales",

        "Summary table with key metrics"

    ]),

    ("Strategic Recommendations", [

        "Actionable insights for improving sales",

        "Focus areas for agents, products, and branches"

    ]),

    ("Thank You / Q&A", [

        "Summary of trends",

        "Invitation for questions"

    ])

]

# Add bullet slides

for title, bullets in slides_content:

    slide = prs.slides.add_slide(bullet_slide_layout)

    shapes = slide.shapes

    title_shape = shapes.title

    body_shape = shapes.placeholders[1]

    title_shape.text = title

    tf = body_shape.text_frame

    tf.clear()

    for bullet in bullets:

        p = tf.add_paragraph()

        p.text = bullet

        p.level = 0

        p.font.size = Pt(18)

        p.font.color.rgb = RGBColor(0, 0, 0)

# Save presentation

output_path = "/mnt/data/Sales_Performance_Report_2015_2016.pptx"

prs.save(output_path)

print(f"Presentation saved to {output_path}")