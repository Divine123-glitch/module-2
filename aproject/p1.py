# generate_sales_ppt.py

# This script generates a PowerPoint presentation (.pptx) based on sales performance data from 2015-2016.

# To run: Save this as generate_sales_ppt.py, install python-pptx with 'pip install python-pptx', then 'python generate_sales_ppt.py'

# Output: Sales_Performance_Report.pptx

from pptx import Presentation

from pptx.util import Inches, Pt

from pptx.enum.text import PP_ALIGN

# Create a new presentation

prs = Presentation()

# Helper function to add a slide with title and bullets

def add_slide(prs, title, bullets, layout=1):  # layout=1 is Title and Content

    slide_layout = prs.slide_layouts[layout]

    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title

    title_shape.text = title

    

    body_shape = slide.placeholders[1]

    tf = body_shape.text_frame

    tf.clear()  # Clear default text

    

    for bullet in bullets:

        p = tf.add_paragraph()

        p.text = bullet

        # Simple indent for sub-bullets (if bullet starts with spaces or tabs)

        if bullet.startswith('  '):

            p.level = 1

        p.alignment = PP_ALIGN.LEFT

        for paragraph in tf.paragraphs:

            for run in paragraph.runs:

                run.font.size = Pt(18)

    

    # Add footer to the slide (using a text box if no placeholder)

    footer_txbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))

    footer_frame = footer_txbox.text_frame

    footer_frame.text = "United to Prevent Violence Against Women and Girls | www.cewhin.com"

    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for run in footer_frame.paragraphs[0].runs:

        run.font.size = Pt(12)

    

    return slide

# Slide 1: Title Slide

title_slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title

title.text = "Sales Performance Report: 2015-2016"

subtitle = slide.placeholders[1]

subtitle.text = "Key Insights from Branch Data"

# Footer

footer_txbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))

footer_frame = footer_txbox.text_frame

footer_frame.text = "United to Prevent Violence Against Women and Girls | www.cewhin.com"

footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

for run in footer_frame.paragraphs[0].runs:

    run.font.size = Pt(12)

# Slide 2: 2015 Agent Performance

add_slide(prs, "2015 Agent Performance", [

    "Tolu had the highest sales revenue",

    "  Total Rev: 2,883,445 / 21.08%",

    "Tonye had the lowest sales revenue",

    "  Total Per: 289.12 / 2.79%"

])

# Slide 3: 2015 Product Performance

add_slide(prs, "2015 Product Performance", [

    "HP had the highest sales revenue",

    "  Total Par: 5,814 / 414 units",

    "Lenovo had the lowest sales revenue",

    "  Total Par: 208.80 / 160 units",

    "Average Price of the year 2015: 200.00"

])

# Slide 4: 2015 Branch Performance

add_slide(prs, "2015 Branch Performance", [

    "Ijoh is the highest performing branch",

    "  Total Rev: 7,305.56 / 70.45%",

    "GRT is the lowest performing branch",

    "  Total Revenue: 808.38 / 7.8%"

])

# Slide 5: 2016 Agent Performance

add_slide(prs, "2016 Agent Performance", [

    "Chinma had the highest sales revenue",

    "  Total Rev: 3,102 / 33.51%",

    "Torbari had the lowest sales revenue",

    "  Total Rev: 57.71 / 0.62%"

])

# Slide 6: 2016 Product Performance

add_slide(prs, "2016 Product Performance", [

    "Apple had the highest performing sales",

    "  Total Per: 3,247 / 308 units",

    "Apple had the lowest (note: possible data overlap; lowest units: 250 / 2 units)",

    "Average price of 2016: 125.00"

])

# Slide 7: 2016 Branch Performance

add_slide(prs, "2016 Branch Performance", [

    "GRA had the highest performance",

    "  Total Rev: 5,193 / 56%",

    "Addah Town had the lowest performance",

    "  Total Rev: 231.12 / 2.5%"

])

# Slide 8: 2016 Revenue by Month

add_slide(prs, "2016 Revenue by Month", [

    "July had the highest performance",

    "  Total Rev: 1,646",

    "March had the lowest performance",

    "  Total Rev: 167.44"

])

# Slide 9: Revenue by Branch (Overall)

add_slide(prs, "Revenue by Branch (Overall)", [

    "Ijoh Branch has the highest sales revenue",

    "  Total Revenue: 11,139.07 / 56.73%",

    "Town Branch has the lowest sales revenue",

    "  Total Revenue: 2,486.42 / 10.67%"

])

# Slide 10: Noticeable Insights

add_slide(prs, "Noticeable Insights", [

    "Only 3 out of 11 agents sold Apple products (Tolu, Emeka, Chinma)",

    "Only 5 out of 11 agents sold company products (Blessing, Ibrahim, Torbari, Chinma, Uche)",

    "Only 5 out of 11 agents sold HP products – DELL (Emeka, China, George, Blessing, Tolu)",

    "All 11 agents sold HP products",

    "Only 8 agents out of 11 sold Lenovo products (Tolu, George, Blessing, Tonye, Uche, Chinedu, Ibrahim, Tunde)",

    "No agent sold all 5 products",

    "Blessing and Uche sold the highest number of distinct products (4)",

    "Torbari and Tunde sold the lowest number of distinct products (2)"

])

# Slide 11: Over the 2 Years – Agent & Product Performance

add_slide(prs, "Over the 2 Years – Agent & Product Performance", [

    "Agent Performance:",

    "  Emeka had the highest sales revenue: 3,109.44 (15.84%)",

    "  Torbari had the lowest sales revenue: 536.75 (2.73%)",

    "Product Performance:",

    "  HP has the highest sales revenue: 955k / 722 units sold",

    "  Apple has the lowest sales revenue: 1.5k / 10 units sold"

])

# Slide 12: Over the 2 Years – Year Performance

add_slide(prs, "Over the 2 Years – Year Performance", [

    "2015 is the highest performing year",

    "  Total Revenue: 10,369.54",

    "  Units Sold: 943",

    "2016 is the lowest performing year",

    "  Total Revenue: 9,258.39",

    "  Units Sold: MK (data unclear; possibly missing)"

])

# Slide 13: Revenue by Month Over the 2 Years

add_slide(prs, "Revenue by Month Over the 2 Years", [

    "December has the highest revenue by month",

    "  We have more sales in December than any other month",

    "March has the least revenue by month",

    "  We have less sales in the month of March"

])

# Slide 14: Thank You / Q&A

add_slide(prs, "Thank You / Q&A", [

    "Summary: Key trends show strong performance in 2015, HP dominance, and seasonal peaks in December/July.",

    "Contact: For more details, visit www.cewhin.com"

])

# Save the presentation

prs.save('Sales_Performance_Report.pptx')

print("PPTX file generated successfully: Sales_Performance_Report.pptx")

